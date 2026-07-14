import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import io
import pandas as pd
import numpy as np

# ----------------------------------------------------------------------
# 🎨 PALETA "DARK ANALYTICS" (MODO OSCURO DE ALTO CONTRASTE)
# ----------------------------------------------------------------------
FONDO_DARK   = RGBColor(11, 19, 43)      # Azul medianoche casi negro
PANEL_DARK   = RGBColor(28, 37, 65)      # Azul grisáceo para tarjetas
CYAN_NEON    = RGBColor(0, 180, 216)     # Acento principal (Títulos y Data)
BLANCO_TXT   = RGBColor(240, 240, 245)   # Blanco roto para lectura
GRIS_TXT     = RGBColor(170, 180, 190)   # Gris para subtítulos
VERDE_NEON   = RGBColor(6, 214, 160)     # Éxito / Positivo
ROJO_NEON    = RGBColor(239, 35, 60)     # Riesgo / Negativo

FONT_FAM = "Segoe UI"

# ----------------------------------------------------------------------
# 🧱 COMPONENTES DE DISEÑO DARK MODE
# ----------------------------------------------------------------------
def crear_encabezado_dark(slide, titulo):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = FONDO_DARK
    
    linea = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.1))
    linea.fill.solid()
    linea.fill.fore_color.rgb = CYAN_NEON
    linea.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = titulo.upper()
    p.font.name = FONT_FAM
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLANCO_TXT

def crear_kpi_dark(slide, x, y, w, h, titulo, valor, color_destaque):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PANEL_DARK
    bg.line.fill.background()
    
    ind = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.1), h)
    ind.fill.solid()
    ind.fill.fore_color.rgb = color_destaque
    ind.line.fill.background()

    tx = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), w - Inches(0.3), h - Inches(0.2))
    tf = tx.text_frame
    
    p1 = tf.paragraphs[0]
    p1.text = titulo
    p1.font.name = FONT_FAM
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = GRIS_TXT
    
    p2 = tf.add_paragraph()
    p2.text = str(valor)
    p2.font.name = FONT_FAM
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = color_destaque

def caja_explicacion_dark(slide, x, y, w, h, titulo, texto_dinamico):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PANEL_DARK
    bg.line.fill.background()
    
    tx = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), h - Inches(0.3))
    tf = tx.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = titulo
    p1.font.name = FONT_FAM
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_NEON
    
    p2 = tf.add_paragraph()
    p2.text = texto_dinamico
    p2.font.name = FONT_FAM
    p2.font.size = Pt(14)
    p2.font.color.rgb = BLANCO_TXT
    p2.space_before = Pt(8)

def configurar_grafico_oscuro(ax, fig):
    fig.patch.set_alpha(0.0)
    ax.patch.set_alpha(0.0)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#888888')
    ax.spines['left'].set_color('#888888')
    ax.tick_params(colors='white')
    ax.xaxis.label.set_color('white')
    ax.yaxis.label.set_color('white')

# ----------------------------------------------------------------------
# 🚀 MOTOR GENERADOR DEL REPORTE (8 DIAPOSITIVAS)
# ----------------------------------------------------------------------
def generar_ppt_ejecutivo(df):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    
    total_resenas = len(df)
    df_pos = df[df["sentiment"] == "positivo"]
    df_neg = df[df["sentiment"] == "negativo"]
    
    pct_pos = round((len(df_pos) / total_resenas) * 100, 1) if total_resenas > 0 else 0
    pct_neg = round((len(df_neg) / total_resenas) * 100, 1) if total_resenas > 0 else 0
    
    unica_categoria = df["product_category"].nunique() == 1 if "product_category" in df.columns else False
    cat_nombre = df["product_category"].iloc[0] if unica_categoria else "Catálogo General"
    
    top_cat_neg = df_neg["product_category"].value_counts().idxmax() if "product_category" in df.columns and not df_neg.empty else "N/A"
    top_topic_neg = df_neg["tematica_resena"].value_counts().idxmax() if "tematica_resena" in df.columns and not df_neg.empty else "N/A"
    top_topic_pos = df_pos["tematica_resena"].value_counts().idxmax() if "tematica_resena" in df.columns and not df_pos.empty else "N/A"

    # ==================================================================
    # SLIDE 1: PORTADA DARK TECH (CON IMAGEN)
    # ==================================================================
    s1 = prs.slides.add_slide(blank)
    s1.background.fill.solid()
    s1.background.fill.fore_color.rgb = FONDO_DARK
    
    # 1. Imagen de fondo en la mitad derecha
    try:
        url_img = "https://images.unsplash.com/photo-1551288049-bebda4e38f71?q=80&w=800&auto=format&fit=crop"
        res = requests.get(url_img)
        s1.shapes.add_picture(io.BytesIO(res.content), Inches(6.0), Inches(0), width=Inches(7.33), height=Inches(7.5))
    except:
        pass 

    # 2. Panel oscuro estricto en la mitad izquierda (Evita desbordes)
    panel = s1.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(6.5), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = FONDO_DARK
    panel.line.fill.background()

    # 3. Textos bien ajustados
    tx = s1.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5), Inches(5.0))
    tf = tx.text_frame
    tf.word_wrap = True 
    
    p = tf.paragraphs[0]
    p.text = "PROYECTO DE TITULACIÓN"
    p.font.name = FONT_FAM
    p.font.size = Pt(14)
    p.font.color.rgb = CYAN_NEON
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = "Diseño e Implementación de una Plataforma Escalable en la Nube para Análisis de Sentimiento"
    p2.font.name = FONT_FAM
    p2.font.size = Pt(32) # Tamaño seguro
    p2.font.color.rgb = BLANCO_TXT
    p2.font.bold = True
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = f"\nMuestra analizada: {cat_nombre}\n{total_resenas} interacciones evaluadas por IA"
    p3.font.name = FONT_FAM
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRIS_TXT

    # ==================================================================
    # SLIDE 2: RESUMEN EJECUTIVO (KPIs)
    # ==================================================================
    s2 = prs.slides.add_slide(blank)
    crear_encabezado_dark(s2, "1. Resumen Ejecutivo de la Auditoría")

    crear_kpi_dark(s2, Inches(0.5), Inches(1.2), Inches(3.0), Inches(1.2), "Total Analizado", f"{total_resenas:,}", CYAN_NEON)
    crear_kpi_dark(s2, Inches(3.8), Inches(1.2), Inches(3.0), Inches(1.2), "Tasa de Aprobación", f"{pct_pos}%", VERDE_NEON)
    crear_kpi_dark(s2, Inches(7.1), Inches(1.2), Inches(3.0), Inches(1.2), "Riesgo de Fuga", f"{pct_neg}%", ROJO_NEON)

    texto_ejecutivo = (
        f"El pipeline de Machine Learning procesó exitosamente {total_resenas} registros correspondientes al segmento '{cat_nombre}'.\n\n"
        f"El desempeño general muestra una tasa de aprobación del {pct_pos}%. El algoritmo de extracción semántica señala que el mayor diferenciador competitivo actualmente es '{top_topic_pos}'.\n\n"
        f"Alerta Operativa: Se detectó un {pct_neg}% de interacciones críticas. El principal destructor de valor extraído por el modelo es '{top_topic_neg}'."
    )
    caja_explicacion_dark(s2, Inches(0.5), Inches(2.8), Inches(12.3), Inches(4.0), "Diagnóstico Generado por IA", texto_ejecutivo)

    # ==================================================================
    # SLIDE 3: DISTRIBUCIÓN DE SENTIMIENTOS
    # ==================================================================
    s3 = prs.slides.add_slide(blank)
    crear_encabezado_dark(s3, "2. Distribución de Sentimiento Global")

    conteo = df["sentiment"].value_counts()
    fig, ax = plt.subplots(figsize=(5, 4.5))
    colores_pie = {"positivo": "#06d6a0", "negativo": "#ef233c", "neutral": "#5c677d"}
    ax.pie(conteo, labels=conteo.index, autopct="%1.1f%%", colors=[colores_pie.get(x, "#ccc") for x in conteo.index], startangle=90, textprops={'weight':'bold', 'color':'white'})
    fig.patch.set_alpha(0.0)
    
    img_d = io.BytesIO()
    plt.savefig(img_d, bbox_inches="tight", transparent=True, dpi=150)
    img_d.seek(0)
    plt.close()
    s3.shapes.add_picture(img_d, Inches(0.5), Inches(1.5), width=Inches(5.0))

    txt_donut = f"El modelo clasificador calculó la polaridad matemática exacta de las {total_resenas} reseñas extraídas.\n\n"
    if pct_neg > 20:
        txt_donut += f"⚠️ RIESGO ALTO: La tasa de negatividad ({pct_neg}%) supera los márgenes operativos saludables. Es imperativo activar las recomendaciones al final de este informe."
    else:
        txt_donut += f"✅ SALUDABLE: La retención proyectada es alta ({pct_pos}%). Los casos negativos están dentro del margen de fricción normal."
    caja_explicacion_dark(s3, Inches(6.0), Inches(2.0), Inches(6.5), Inches(4.0), "Evaluación de Polaridad", txt_donut)

    # ==================================================================
    # SLIDE 4: CATEGORÍAS CON MAYOR IMPACTO
    # ==================================================================
    if "product_category" in df.columns:
        s4 = prs.slides.add_slide(blank)
        crear_encabezado_dark(s4, "3. Interacción Comercial por Categorías")

        categorias = df["product_category"].value_counts().head(5)
        fig, ax = plt.subplots(figsize=(6, 4))
        bars = ax.barh(categorias.index, categorias.values, color="#00b4d8", height=0.6)
        configurar_grafico_oscuro(ax, fig)
        ax.invert_yaxis()
        
        for bar in bars:
            ax.text(bar.get_width()+0.5, bar.get_y()+bar.get_height()/2, str(int(bar.get_width())), va='center', color='white', weight='bold')
            
        img_c = io.BytesIO()
        plt.tight_layout()
        plt.savefig(img_c, bbox_inches="tight", transparent=True, dpi=150)
        img_c.seek(0)
        plt.close()
        s4.shapes.add_picture(img_c, Inches(0.5), Inches(1.8), width=Inches(6.5))

        txt_cat = f"Al analizar la estructura de los datos, el sistema detecta que el segmento con mayor tráfico de opiniones es '{categorias.idxmax()}' con {categorias.max()} reseñas procesadas.\n\nEsto define el peso relativo de las decisiones: los presupuestos de optimización deben enfocarse primordialmente en esta categoría de alta demanda."
        caja_explicacion_dark(s4, Inches(7.3), Inches(2.0), Inches(5.5), Inches(4.0), "Tráfico de Reseñas", txt_cat)

    # ==================================================================
    # SLIDE 5: FORTALEZAS (LO QUE DEBEMOS POTENCIAR)
    # ==================================================================
    if "tematica_resena" in df.columns and not df_pos.empty:
        s5 = prs.slides.add_slide(blank)
        crear_encabezado_dark(s5, "4. Pilares de Éxito y Fortalezas (BERTopic)")

        tops_p = df_pos["tematica_resena"].value_counts().head(4)
        fig, ax = plt.subplots(figsize=(6, 4))
        bars = ax.barh(tops_p.index, tops_p.values, color="#06d6a0", height=0.6)
        configurar_grafico_oscuro(ax, fig)
        ax.invert_yaxis()
        
        for bar in bars:
            ax.text(bar.get_width()+0.5, bar.get_y()+bar.get_height()/2, str(int(bar.get_width())), va='center', color='white', weight='bold')
            
        img_tp = io.BytesIO()
        plt.tight_layout()
        plt.savefig(img_tp, bbox_inches="tight", transparent=True, dpi=150)
        img_tp.seek(0)
        plt.close()
        s5.shapes.add_picture(img_tp, Inches(0.5), Inches(1.8), width=Inches(6.5))

        txt_fortaleza = f"El modelo extrajo los tópicos exactos dentro de las opiniones positivas. El atributo más elogiado es '{tops_p.idxmax()}' ({tops_p.max()} menciones).\n\nEstrategia: El equipo de Marketing debe utilizar estos tópicos como los 'ganchos' principales en futuras campañas publicitarias, ya que representan el valor real percibido por los clientes."
        caja_explicacion_dark(s5, Inches(7.3), Inches(2.0), Inches(5.5), Inches(4.0), "Oportunidades de Crecimiento", txt_fortaleza)

    # ==================================================================
    # SLIDE 6: PROBLEMAS ENCONTRADOS (CAUSA RAÍZ)
    # ==================================================================
    if "tematica_resena" in df.columns and not df_neg.empty:
        s6 = prs.slides.add_slide(blank)
        crear_encabezado_dark(s6, "5. Focos Críticos y Fallas (Causa Raíz)")

        tops_n = df_neg["tematica_resena"].value_counts().head(4)
        fig, ax = plt.subplots(figsize=(6, 4))
        bars = ax.barh(tops_n.index, tops_n.values, color="#ef233c", height=0.6)
        configurar_grafico_oscuro(ax, fig)
        ax.invert_yaxis()
        
        for bar in bars:
            ax.text(bar.get_width()+0.5, bar.get_y()+bar.get_height()/2, str(int(bar.get_width())), va='center', color='white', weight='bold')
            
        img_tn = io.BytesIO()
        plt.tight_layout()
        plt.savefig(img_tn, bbox_inches="tight", transparent=True, dpi=150)
        img_tn.seek(0)
        plt.close()
        s6.shapes.add_picture(img_tn, Inches(0.5), Inches(1.8), width=Inches(6.5))

        txt_causa = f"El algoritmo agrupó semánticamente las quejas para descubrir por qué los clientes se van. La falla operativa número uno es '{tops_n.idxmax()}', con {tops_n.max()} incidentes confirmados.\n\nOperaciones debe destinar recursos inmediatos para solucionar exclusivamente esta variable, lo que reducirá la tasa de devolución de forma escalable."
        caja_explicacion_dark(s6, Inches(7.3), Inches(2.0), Inches(5.5), Inches(4.0), "Análisis de Fricción Operativa", txt_causa)

    # ==================================================================
    # SLIDE 7: COMENTARIOS DESTACADOS
    # ==================================================================
    s7 = prs.slides.add_slide(blank)
    crear_encabezado_dark(s7, "6. La Voz del Mercado (Verbatims Destacados)")

    y_pos = 1.2
    for sent, titulo, color in [("positivo", "Promotor (Muestra de Lealtad y Satisfacción)", VERDE_NEON), 
                                ("negativo", f"Detractor Crítico (Riesgo por '{top_topic_neg}')", ROJO_NEON)]:
        sub = df[df.sentiment == sent]
        texto = sub.iloc[0]["review"][:250] + "..." if not sub.empty else "No hay registros disponibles."
        
        caja_explicacion_dark(s7, Inches(0.5), Inches(y_pos), Inches(12.33), Inches(1.8), titulo, f'"{texto}"')
        s7.shapes[-1].text_frame.paragraphs[0].font.color.rgb = color
        
        y_pos += 2.2

    # ==================================================================
    # SLIDE 8: RECOMENDACIONES Y PLAN DE ACCIÓN
    # ==================================================================
    s8 = prs.slides.add_slide(blank)
    crear_encabezado_dark(s8, "7. Plan de Acción y Conclusión")

    tabla = s8.shapes.add_table(4, 3, Inches(0.5), Inches(1.2), Inches(12.33), Inches(2.8)).table
    tabla.columns[0].width = Inches(1.8)
    tabla.columns[1].width = Inches(7.53)
    tabla.columns[2].width = Inches(3.0)

    for i, head in enumerate(["NIVEL DE RIESGO", "ESTRATEGIA PRESCRIPTIVA SUGERIDA", "ÁREA DE EJECUCIÓN"]):
        celda = tabla.cell(0, i)
        celda.fill.solid()
        celda.fill.fore_color.rgb = CYAN_NEON
        p = celda.text_frame.paragraphs[0]
        p.text = head
        p.font.color.rgb = FONDO_DARK
        p.font.bold = True
        p.font.size = Pt(12)

    accion_critica = f"Auditar con urgencia los procesos logísticos/técnicos ligados a '{top_topic_neg}' para frenar el {pct_neg}% de fuga."
    accion_marketing = f"Inyectar las palabras clave asociadas a '{top_topic_pos}' en todas las nuevas campañas (Ads y Redes)."

    filas = [
        ("CRÍTICO", accion_critica, "Operaciones y Calidad"),
        ("ALTO", "Contactar proactivamente a los clientes insatisfechos para evitar viralidad negativa online.", "Servicio al Cliente"),
        ("OPORTUNIDAD", accion_marketing, "Marketing y Growth")
    ]

    for i, (prio, accion, area) in enumerate(filas):
        tabla.cell(i+1, 0).text = prio
        tabla.cell(i+1, 1).text = accion
        tabla.cell(i+1, 2).text = area
        
        for j in range(3):
            celda = tabla.cell(i+1, j)
            celda.fill.solid()
            celda.fill.fore_color.rgb = PANEL_DARK
            pf = celda.text_frame.paragraphs[0]
            pf.font.size = Pt(13)
            pf.font.color.rgb = BLANCO_TXT
            if prio == "CRÍTICO" and j == 0:
                pf.font.color.rgb = ROJO_NEON
                pf.font.bold = True
            elif prio == "OPORTUNIDAD" and j == 0:
                pf.font.color.rgb = VERDE_NEON
                pf.font.bold = True

    txt_cierre = "Conclusión: La implementación de esta plataforma demuestra que es viable integrar Machine Learning en la nube para procesar reseñas masivas. Permite aislar anomalías en tiempo real, transformar texto ambiguo en datos matemáticos estructurados y ahorrar cientos de horas operativas en análisis manual, dotando a la gerencia de inteligencia de negocios aplicable al instante."
    caja_explicacion_dark(s8, Inches(0.5), Inches(4.5), Inches(12.33), Inches(2.2), "Cierre de Auditoría del Proyecto", txt_cierre)

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream